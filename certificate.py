"""
==============================================================================
 Accredify — Certificate Automation Suite
 (Enterprise Dashboard UI — Fluent / Azure Portal / GitHub Desktop inspired)
==============================================================================
This file is a UI/UX redesign of the original Certificate Automation Suite.
ALL backend logic — Excel parsing, PPTX placeholder replacement, offline QR
payload generation, PowerPoint COM export to PDF, Outlook COM sending, and
Safe-Mode throttling — is IDENTICAL to the original implementation. Only the
following were added, and only as additive, optional hooks (default = None,
so nothing breaks if unused):
    - `stage_callback` and `row_callback` parameters on `process_roster()`,
      used purely to power the new live execution dashboard (current
      pipeline stage, ETA, running success/fail counters). The actual
      processing steps, their order, and their behavior are unchanged.
    - `generate_single_certificate()`, a new helper that reuses the exact
      same personalization / QR / COM functions to render a preview for
      ONE attendee on demand (used by the "Preview Certificate" action).

------------------------------------------------------------------------
SETUP INSTRUCTIONS
------------------------------------------------------------------------
1. Windows 10/11, Microsoft Office (classic PowerPoint + classic Outlook —
   NOT the new standalone Outlook app, which does not support COM
   automation), Python 3.9+.

2. Create and activate a virtual environment:
       python -m venv venv
       venv\\Scripts\\activate

3. Install dependencies (Streamlit 1.36+ required for native Material
   icons on buttons; 1.37+ required for st.dialog modals):
       pip install "streamlit>=1.37" pandas openpyxl python-pptx qrcode[pil] Pillow pywin32

   requirements.txt
   ---------------------------------
   streamlit>=1.37
   pandas>=2.0
   openpyxl>=3.1
   python-pptx>=0.6.23
   qrcode[pil]>=7.4
   Pillow>=10.0
   pywin32>=306
   ---------------------------------

4. Run pywin32's post-install script once:
       python venv\\Scripts\\pywin32_postinstall.py -install

5. Open classic Outlook desktop once, fully signed in, BEFORE launching
   this app.

6. Launch:
       streamlit run ieee_vit_certificate_suite.py

------------------------------------------------------------------------
HONEST UI LIMITATIONS (Streamlit constraints — flagged rather than faked)
------------------------------------------------------------------------
- The bottom action bar is styled to look docked, but Streamlit does not
  support true OS-level fixed/sticky footers independent of scroll
  position with full reliability across all browsers; it uses CSS
  `position: sticky` on its wrapping block, which docks it within the
  main content column as you scroll, but will not float over the
  sidebar or persist if Streamlit's own layout re-flows on some browsers.
- There is no supported public API to switch the active `st.tabs()` tab
  from Python code, so the "Preview Template" button on the dashboard
  opens a modal dialog with the rendered slide instead of jumping to the
  Layout Preview tab.
- "Sortable columns" / "sticky headers" on the attendee table are
  implemented manually (a fixed header row above a scrollable, fixed-
  height container + a sort dropdown), since `st.dataframe` does not yet
  support inline per-row action buttons.
- "Zoom in / out" on the layout preview is implemented via a display-
  width slider on the rendered PNG image, not true vector/DPI zoom.
==============================================================================
"""

import streamlit as st
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Emu
import qrcode
from PIL import Image
import os
import io
import time
import random
import tempfile
import shutil
import traceback
import platform
from datetime import datetime, date

IS_WINDOWS = platform.system() == "Windows"
WINDOWS_COM_AVAILABLE = False
COM_IMPORT_ERROR = None

if IS_WINDOWS:
    try:
        import win32com.client
        import pythoncom
        WINDOWS_COM_AVAILABLE = True
    except ImportError as exc:
        COM_IMPORT_ERROR = str(exc)
else:
    COM_IMPORT_ERROR = "Not running on Windows — COM automation (PowerPoint/Outlook) is unavailable."

PP_SAVE_AS_PDF = 32
OL_MAIL_ITEM = 0


# ==========================================================================
# PAGE CONFIGURATION
# ==========================================================================
st.set_page_config(
    page_title="Accredify | Certificate Automation Suite",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ==========================================================================
# CUSTOM CSS — Fluent / Azure Portal inspired enterprise dark theme
# ==========================================================================
CUSTOM_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Hanken+Grotesk:wght@400;500;600;700;800&family=JetBrains+Mono:wght@400;500;600&display=swap');
:root {
    /* ---- Elevation layer 0: page background ---- */
    --bg: #0B1120;
    /* ---- Elevation layer 1: static cards/panels ---- */
    --card: #151D2E;
    /* ---- Elevation layer 2: interactive/secondary surface ---- */
    --surface: #1B2436;
    --surface-hover: #232D45;
    --border: #2A3450;
    --accent: #6366F1;
    --accent-2: #6366F1;
    --accent-hover: #818CF8;
    --success: #22C55E;
    --warning: #F59E0B;
    --error: #EF4444;
    --text: #E5E9F5;
    --text-dim: #8A93A8;
    --radius: 12px;
    --radius-sm: 10px;
    /* ---- 8px spacing scale ---- */
    --sp-1: 8px;
    --sp-2: 16px;
    --sp-3: 24px;
    --sp-4: 32px;
}

html, body, [data-testid="stAppViewContainer"] {
    background: var(--bg) !important;
    color: var(--text) !important;
    font-family: "Hanken Grotesk", "Segoe UI", "Inter", -apple-system, sans-serif;
}

[data-testid="stHeader"] { background: transparent !important; }

[data-testid="stSidebar"] {
    background: #0E1524 !important;
    border-right: 1px solid var(--border);
    width: 240px !important;
}
[data-testid="stSidebar"] * { color: var(--text) !important; }
[data-testid="stSidebar"] .block-container { padding-top: var(--sp-2); }

h1, h2, h3, h4 { color: var(--text) !important; font-weight: 700 !important; letter-spacing: -0.01em; }
p, span, label, div { color: var(--text); }

/* ---------- Reduce default block gaps slightly for tighter rhythm ---------- */
.block-container { padding-top: var(--sp-3) !important; }
div[data-testid="stVerticalBlock"] > div { gap: var(--sp-1); }

/* ---------- Top nav bar ---------- */
.topnav {
    display: flex; align-items: center; justify-content: space-between;
    padding: var(--sp-2) var(--sp-3); background: var(--card); border: 1px solid var(--border);
    border-radius: var(--radius); margin-bottom: var(--sp-2);
    box-shadow: 0 1px 2px rgba(0,0,0,0.25);
}
.topnav-brand { display: flex; align-items: center; gap: var(--sp-1); font-weight: 800; font-size: 1rem; color: var(--accent); }
.topnav-title { font-weight: 700; font-size: 1.15rem; color: var(--text); }
.topnav-status { display: flex; gap: var(--sp-1); flex-wrap: wrap; }

/* ---------- Status pills ---------- */
.pill {
    display: inline-flex; align-items: center; gap: 6px;
    padding: 4px 12px; border-radius: 999px; font-size: 0.74rem; font-weight: 600;
    border: 1px solid var(--border); white-space: nowrap;
}
.pill-ok { background: rgba(34,197,94,0.14); color: #4ADE80; border-color: rgba(34,197,94,0.35); }
.pill-bad { background: rgba(239,68,68,0.14); color: #F87171; border-color: rgba(239,68,68,0.35); }
.pill-warn { background: rgba(245,158,11,0.16); color: #FBBF24; border-color: rgba(245,158,11,0.35); }
.pill-neutral { background: rgba(99,102,241,0.14); color: #A5B4FC; border-color: rgba(99,102,241,0.32); }
.pill-waiting { background: var(--surface); color: var(--text-dim); border-color: var(--border); }

/* ---------- Workflow wizard (Fluent-style step indicators) ---------- */
.wizard { display: flex; align-items: center; gap: 4px; margin: 0 0 var(--sp-3) 0; }
.wizard-step {
    flex: 1; display: flex; align-items: center; justify-content: center; gap: 6px;
    text-align: center; padding: 10px 6px; border-radius: var(--radius-sm);
    background: var(--card); border: 1px solid var(--border); color: var(--text-dim);
    font-size: 0.78rem; font-weight: 600; transition: all 0.2s ease;
}
.wizard-step .step-dot {
    width: 18px; height: 18px; border-radius: 50%; display: inline-flex; align-items: center;
    justify-content: center; font-size: 0.68rem; font-weight: 800; flex-shrink: 0;
    background: var(--surface); color: var(--text-dim); border: 1px solid var(--border);
}
.wizard-step.active {
    background: rgba(99,102,241,0.14);
    border-color: var(--accent); color: #A5B4FC;
    box-shadow: 0 0 0 1px var(--accent) inset;
}
.wizard-step.active .step-dot { background: var(--accent); color: #fff; border-color: var(--accent); }
.wizard-step.done { border-color: var(--success); color: #4ADE80; background: rgba(34,197,94,0.08); }
.wizard-step.done .step-dot { background: var(--success); color: #06210F; border-color: var(--success); }
.wizard-arrow { color: var(--border); font-size: 0.9rem; }

/* ---------- Generic card (elevation layer 1) ---------- */
.card {
    background: var(--card); border: 1px solid var(--border); border-radius: var(--radius);
    padding: var(--sp-3); margin-bottom: var(--sp-2); transition: border-color 0.15s ease, box-shadow 0.15s ease;
    box-shadow: 0 1px 2px rgba(0,0,0,0.2);
}
.card:hover { border-color: #3A4568; }
.card-title { font-size: 1rem; font-weight: 700; margin-bottom: var(--sp-2); display: flex; align-items: center; gap: var(--sp-1); color: var(--text); }
.card-subtle { color: var(--text-dim); font-size: 0.82rem; margin-top: -6px; margin-bottom: var(--sp-1); }

/* ---------- Primary focal-point cards (Event Details / Upload / Attendee Mgmt / Execute) ---------- */
.card-primary {
    background: var(--card); border: 1px solid var(--border); border-radius: var(--radius);
    padding: var(--sp-4) var(--sp-3); margin-bottom: var(--sp-3);
    border-top: 3px solid var(--accent);
    box-shadow: 0 1px 3px rgba(0,0,0,0.25);
}
.card-primary .card-title { font-size: 1.1rem; }

/* ---------- KPI grid: responsive 2x2 → 4-across, min 220px, no vertical text wrap ---------- */
.kpi-grid {
    display: grid; grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
    gap: var(--sp-2); margin-bottom: var(--sp-2);
}
.kpi-card {
    background: var(--card);
    border: 1px solid var(--border); border-radius: var(--radius);
    padding: var(--sp-2); min-width: 220px; max-width: 260px;
    display: flex; flex-direction: column; gap: 6px;
    box-shadow: 0 1px 2px rgba(0,0,0,0.2);
}
.kpi-top { display: flex; align-items: center; justify-content: space-between; }
.kpi-icon { font-size: 1.3rem; }
.kpi-label { color: var(--text-dim); font-size: 0.72rem; font-weight: 700; text-transform: uppercase; letter-spacing: 0.03em; white-space: nowrap; }
.kpi-value { font-size: 1.35rem; font-weight: 800; white-space: nowrap; color: var(--text); }
.kpi-secondary { color: var(--text-dim); font-size: 0.74rem; }

/* ---------- Placeholder / validation tiles: 3-state (waiting / valid / error) ---------- */
.tile { border-radius: var(--radius-sm); border: 1px solid var(--border); padding: var(--sp-1) var(--sp-2); text-align: center; font-weight: 700; font-size: 0.85rem; }
.tile-ok { background: rgba(34,197,94,0.12); border-color: rgba(34,197,94,0.4); color: #4ADE80; }
.tile-bad { background: rgba(239,68,68,0.12); border-color: rgba(239,68,68,0.4); color: #F87171; }
.tile-waiting { background: var(--surface); border-color: var(--border); color: var(--text-dim); }

/* ---------- Terminal console (professional, with log levels) ---------- */
.terminal {
    background: #060A14; border: 1px solid var(--border); border-radius: var(--radius-sm);
    padding: var(--sp-2); font-family: "JetBrains Mono", "Cascadia Code", "Consolas", monospace; font-size: 0.82rem;
    max-height: 340px; overflow-y: auto; color: #C7D2E0; line-height: 1.65;
}
.terminal .ts { color: #5B6B85; }
.terminal .lvl { font-weight: 800; padding: 0 2px; }
.terminal .lvl-info { color: #38BDF8; }
.terminal .lvl-success { color: #4ADE80; }
.terminal .lvl-warning { color: #FBBF24; }
.terminal .lvl-error { color: #F87171; }

/* ---------- Attendee table ---------- */
.tbl-header {
    display: grid; grid-template-columns: 2fr 2.2fr 1.1fr 2fr;
    gap: var(--sp-1); padding: var(--sp-1) var(--sp-2); background: var(--surface);
    border: 1px solid var(--border); position: sticky; top: 0; z-index: 5;
    border-radius: var(--radius-sm) var(--radius-sm) 0 0; font-size: 0.72rem;
    font-weight: 700; color: var(--text-dim); text-transform: uppercase; letter-spacing: 0.03em;
}
.tbl-row-wrap { border: 1px solid var(--border); border-top: none; border-radius: 0 0 var(--radius-sm) var(--radius-sm); background: var(--card); }
.tbl-pagination { display: flex; align-items: center; justify-content: space-between; margin-top: var(--sp-1); color: var(--text-dim); font-size: 0.8rem; }

/* ---------- Upload cards ---------- */
.upload-illustration {
    border: 1.5px dashed var(--border); border-radius: var(--radius-sm); background: var(--surface);
    padding: var(--sp-2); text-align: center; margin-bottom: 6px;
}
.upload-illustration .upicon { font-size: 1.6rem; margin-bottom: 2px; }
.upload-meta { color: var(--text-dim); font-size: 0.74rem; margin-top: 2px; }
.upload-success { background: rgba(34,197,94,0.12); border: 1px solid rgba(34,197,94,0.35); border-radius: var(--radius-sm); padding: var(--sp-1) var(--sp-2); margin-top: 6px; font-size: 0.8rem; color: #4ADE80; }

/* ---------- Inputs (elevation layer 2 — interactive surface) ---------- */
.stTextInput input, .stDateInput input, .stTextArea textarea, .stNumberInput input, .stSelectbox > div {
    background-color: var(--surface) !important; color: var(--text) !important;
    border-radius: var(--radius-sm) !important; border: 1px solid var(--border) !important;
}
.stTextInput input:focus, .stTextArea textarea:focus { background-color: var(--surface-hover) !important; border-color: var(--accent) !important; }

/* ---------- Buttons: single primary CTA + secondary/outline hierarchy ---------- */
.stButton > button, .stDownloadButton > button {
    background: var(--accent) !important;
    color: #FFFFFF !important; font-weight: 600 !important; border: none !important;
    border-radius: var(--radius-sm) !important; padding: 0.55rem 1.2rem !important;
    box-shadow: 0 1px 2px rgba(99,102,241,0.35); transition: background 0.15s ease, box-shadow 0.15s ease;
}
.stButton > button:hover, .stDownloadButton > button:hover {
    background: var(--accent-hover) !important; box-shadow: 0 2px 8px rgba(99,102,241,0.5);
}
.stButton button[kind="secondary"] {
    background: var(--surface) !important; color: var(--text) !important;
    border: 1px solid var(--border) !important; box-shadow: none; font-weight: 600 !important;
}
.stButton button[kind="secondary"]:hover { background: var(--surface-hover) !important; border-color: var(--accent) !important; box-shadow: none; }

/* ---------- Tabs ---------- */
[data-baseweb="tab-list"] { gap: 6px; }
[data-baseweb="tab"] {
    background: var(--card); border-radius: 10px 10px 0 0; border: 1px solid var(--border);
    color: var(--text-dim); padding: var(--sp-1) var(--sp-2);
}
[aria-selected="true"] { color: #A5B4FC !important; background: var(--surface-hover) !important; border-bottom: 2px solid var(--accent) !important; }

/* ---------- Progress bar ---------- */
.stProgress > div > div { background: var(--accent) !important; }

/* ---------- Sticky-ish action bar ---------- */
.action-bar {
    position: sticky; bottom: 0; z-index: 50;
    background: rgba(11,17,32,0.92); backdrop-filter: blur(6px);
    border: 1px solid var(--border); border-radius: var(--radius);
    padding: var(--sp-2) var(--sp-3); margin-top: var(--sp-1);
    box-shadow: 0 -1px 4px rgba(0,0,0,0.3);
}

/* ---------- File upload card ---------- */
[data-testid="stFileUploader"] {
    background: var(--surface); border: 1.5px dashed var(--border); border-radius: var(--radius-sm); padding: var(--sp-1);
}

hr { border-color: var(--border) !important; }

/* ---------- Native Streamlit widgets — keep them in the dark theme too ---------- */
[data-testid="stExpander"] {
    background: var(--card) !important; border: 1px solid var(--border) !important;
    border-radius: var(--radius-sm) !important;
}
[data-testid="stExpander"] summary { color: var(--text) !important; }

[data-testid="stMetric"] {
    background: var(--surface); border: 1px solid var(--border); border-radius: var(--radius-sm);
    padding: var(--sp-1) var(--sp-2);
}
[data-testid="stMetricLabel"] { color: var(--text-dim) !important; }
[data-testid="stMetricValue"] { color: var(--text) !important; }

[data-testid="stAlert"], .stAlert { background: var(--surface) !important; border: 1px solid var(--border) !important; color: var(--text) !important; }

[data-testid="stDataFrame"], [data-testid="stTable"] { background: var(--card); border-radius: var(--radius-sm); }

/* Checkboxes / radios / sliders */
[data-testid="stCheckbox"] label, [data-testid="stRadio"] label { color: var(--text) !important; }
.stSlider [data-baseweb="slider"] { background: transparent; }

/* Dialog modal (st.dialog) */
[data-testid="stDialog"] > div, div[role="dialog"] {
    background: var(--card) !important; border: 1px solid var(--border) !important; color: var(--text) !important;
}

/* Dropdown / calendar popovers (selectbox, date_input) render in a portal, so target BaseWeb data-attrs directly */
[data-baseweb="popover"] [data-baseweb="menu"], [data-baseweb="popover"] ul {
    background: var(--card) !important; border: 1px solid var(--border) !important;
}
[data-baseweb="menu"] li, [data-baseweb="menu"] li span { color: var(--text) !important; }
[data-baseweb="calendar"] { background: var(--card) !important; color: var(--text) !important; }
[data-baseweb="calendar"] button { color: var(--text) !important; }

/* File uploader drop-zone text */
[data-testid="stFileUploaderDropzone"] { background: var(--surface) !important; }
[data-testid="stFileUploaderDropzoneInstructions"] span, [data-testid="stFileUploaderDropzoneInstructions"] small { color: var(--text-dim) !important; }
</style>
"""
st.markdown(CUSTOM_CSS, unsafe_allow_html=True)


# ==========================================================================
# SESSION STATE INITIALIZATION
# ==========================================================================
_defaults = {
    "run_results_df": None,
    "run_completed": False,
    "roster_df": None,
    "roster_bytes_name": None,
    "template_bytes": None,
    "template_name": None,
    "show_cert_preview": False,
    "preview_row_idx": None,
    "show_template_preview": False,
    "template_preview_zoom": 100,
    "attendee_search": "",
    "attendee_sort": "Row order",
    "attendee_status_filter": "All",
    "run_in_progress": False,
    "active_nav": "dashboard",
}
for _k, _v in _defaults.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v


# ==========================================================================
# BACKEND — CORE LOGIC (unchanged processing behavior)
# ==========================================================================
def generate_unique_id(row_index: int) -> str:
    return f"IEEE-VIT-2026-{row_index + 1:04d}"


def build_offline_verification_payload(unique_id: str, full_name: str,
                                        event_name: str, event_date_str: str) -> str:
    payload = (
        "--- Accredify CERTIFICATE VERIFICATION ---\n"
        "Status: Authentic Record\n"
        f"ID: {unique_id}\n"
        f"Issued To: {full_name}\n"
        f"Event: {event_name}\n"
        f"Date: {event_date_str}\n"
        "-----------------------------------------"
    )
    return payload


def generate_qr_image(payload_text: str, box_size: int = 8, border: int = 2) -> Image.Image:
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=box_size,
        border=border,
    )
    qr.add_data(payload_text)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white").convert("RGB")
    return img


def _replace_in_text_frame(text_frame, mapping: dict) -> None:
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs:
            continue
        full_text = "".join(run.text for run in paragraph.runs)
        replaced_text = full_text
        changed = False
        for placeholder, value in mapping.items():
            if placeholder in replaced_text:
                replaced_text = replaced_text.replace(placeholder, value)
                changed = True
        if changed:
            paragraph.runs[0].text = replaced_text
            for extra_run in paragraph.runs[1:]:
                extra_run.text = ""


def replace_placeholders_in_pptx(prs: Presentation, mapping: dict) -> None:
    for slide in prs.slides:
        _replace_placeholders_in_shapes(slide.shapes, mapping)


def _replace_placeholders_in_shapes(shapes, mapping: dict) -> None:
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            _replace_placeholders_in_shapes(shape.shapes, mapping)
            continue
        if shape.has_text_frame:
            _replace_in_text_frame(shape.text_frame, mapping)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _replace_in_text_frame(cell.text_frame, mapping)


def insert_qr_bottom_right(slide, qr_image_path: str, slide_width_emu: int,
                            slide_height_emu: int, qr_inches: float = 1.3,
                            margin_inches: float = 0.3) -> None:
    qr_size = Inches(qr_inches)
    margin = Inches(margin_inches)
    left = Emu(int(slide_width_emu - qr_size - margin))
    top = Emu(int(slide_height_emu - qr_size - margin))
    slide.shapes.add_picture(qr_image_path, left, top, width=qr_size, height=qr_size)


def pptx_has_all_placeholders(pptx_bytes: bytes) -> dict:
    found = {"{{Name}}": False, "{{Event}}": False, "{{Date}}": False}
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    for key in found:
                        if key in text:
                            found[key] = True
    except Exception:
        pass
    return found


def export_pptx_to_pdf_via_powerpoint(pptx_abs_path: str, pdf_abs_path: str) -> None:
    if not WINDOWS_COM_AVAILABLE:
        raise RuntimeError(
            "PowerPoint COM automation unavailable. This feature requires "
            "Windows with Microsoft PowerPoint installed and pywin32."
        )
    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        # Dispatch (not DispatchEx) is used deliberately — DispatchEx's
        # CoCreateInstanceEx activation path raises "Invalid class string"
        # (-2147221005) on some Windows/Office/pywin32 combinations even
        # though the ProgID is valid and Dispatch resolves it fine.
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            pptx_abs_path, ReadOnly=False, Untitled=False, WithWindow=False
        )
        presentation.SaveAs(pdf_abs_path, PP_SAVE_AS_PDF)
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def render_slide_to_png(pptx_abs_path: str, output_png_abs_path: str,
                         slide_index: int = 1, width_px: int = 1280,
                         height_px: int = 720) -> None:
    if not WINDOWS_COM_AVAILABLE:
        raise RuntimeError(
            "PowerPoint COM automation unavailable. This feature requires "
            "Windows with Microsoft PowerPoint installed and pywin32."
        )
    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            pptx_abs_path, ReadOnly=True, Untitled=False, WithWindow=False
        )
        slide = presentation.Slides(slide_index)
        slide.Export(output_png_abs_path, "PNG", width_px, height_px)
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def send_certificate_via_outlook(to_email: str, subject: str, body: str,
                                  attachment_abs_path: str) -> None:
    if not WINDOWS_COM_AVAILABLE:
        raise RuntimeError(
            "Outlook COM automation unavailable. This feature requires "
            "Windows with classic Outlook installed, pywin32, and an "
            "already logged-in Outlook desktop profile."
        )
    pythoncom.CoInitialize()
    try:
        outlook = win32com.client.Dispatch("Outlook.Application")
        mail = outlook.CreateItem(OL_MAIL_ITEM)
        mail.To = to_email
        mail.Subject = subject
        mail.Body = body
        mail.Attachments.Add(attachment_abs_path)
        mail.Send()
    finally:
        pythoncom.CoUninitialize()


def process_roster(
    df: pd.DataFrame,
    event_name: str,
    event_date_str: str,
    pptx_template_bytes: bytes,
    min_delay: float,
    max_delay: float,
    log_placeholder,
    progress_bar,
    email_subject_template: str,
    email_body_template: str,
    dry_run: bool,
    stage_callback=None,   # NEW, optional, additive: stage_callback(stage_name, idx, total, full_name)
    row_callback=None,     # NEW, optional, additive: row_callback(status)
) -> pd.DataFrame:
    """
    Executes the full end-to-end pipeline for every row in the roster.
    Processing order and behavior are IDENTICAL to the original
    implementation; stage_callback/row_callback are purely observational
    hooks added to power the new live execution dashboard.
    """
    work_dir = tempfile.mkdtemp(prefix="ieee_vit_cert_")
    results = []
    log_lines = []
    total = len(df)

    def _stage(name, idx, full_name):
        if stage_callback:
            try:
                stage_callback(name, idx, total, full_name)
            except Exception:
                pass

    template_path = os.path.abspath(os.path.join(work_dir, "_template.pptx"))
    with open(template_path, "wb") as f:
        f.write(pptx_template_bytes)

    def flush_log():
        log_placeholder.markdown(
            f'<div class="terminal">{"<br>".join(log_lines[-20:])}</div>',
            unsafe_allow_html=True,
        )

    mode_label_html = '<span class="warn">DRY RUN</span>' if dry_run else '<span class="ok">LIVE SEND via Outlook</span>'
    log_lines.append(f'<span class="ts">[{datetime.now().strftime("%H:%M:%S")}]</span> === Pipeline started — Event: {event_name} | Date: {event_date_str} | Attendees: {total} ===')
    log_lines.append(f'<span class="ts">[{datetime.now().strftime("%H:%M:%S")}]</span> Mode: {mode_label_html}')
    flush_log()

    try:
        for idx, row in df.iterrows():
            first_name = str(row.get("First Name", "")).strip()
            last_name = str(row.get("Last Name", "")).strip()
            email = str(row.get("Email ID", "")).strip()
            full_name = f"{first_name} {last_name}".strip()
            unique_id = generate_unique_id(idx)
            status = ""
            ts = datetime.now().strftime("%H:%M:%S")

            log_lines.append(f'<span class="ts">[{ts}]</span> Processing [{idx + 1}/{total}]: {full_name or "UNKNOWN"} ...')
            flush_log()

            try:
                if not first_name and not last_name:
                    raise ValueError("Missing First Name / Last Name in roster row.")
                if not email or "@" not in email or "." not in email.split("@")[-1]:
                    raise ValueError(f"Invalid or missing email address: '{email}'")

                _stage("Generating PPTX", idx, full_name)
                prs = Presentation(template_path)

                _stage("Replacing Placeholders", idx, full_name)
                mapping = {"{{Name}}": full_name, "{{Event}}": event_name, "{{Date}}": event_date_str}
                replace_placeholders_in_pptx(prs, mapping)

                _stage("Embedding QR", idx, full_name)
                payload = build_offline_verification_payload(unique_id, full_name, event_name, event_date_str)
                qr_img = generate_qr_image(payload)
                qr_path = os.path.abspath(os.path.join(work_dir, f"qr_{unique_id}.png"))
                qr_img.save(qr_path)
                target_slide = prs.slides[0]
                insert_qr_bottom_right(target_slide, qr_path, prs.slide_width, prs.slide_height)

                safe_name = (full_name.replace(" ", "_") or f"attendee_{idx + 1}")
                safe_name = "".join(c for c in safe_name if c.isalnum() or c in ("_", "-"))
                individual_pptx_path = os.path.abspath(os.path.join(work_dir, f"{safe_name}_{unique_id}.pptx"))
                prs.save(individual_pptx_path)

                _stage("Exporting PDF", idx, full_name)
                individual_pdf_path = os.path.abspath(os.path.join(work_dir, f"{safe_name}_{unique_id}.pdf"))
                export_pptx_to_pdf_via_powerpoint(individual_pptx_path, individual_pdf_path)

                if not dry_run:
                    _stage("Sending Email", idx, full_name)
                    subject = email_subject_template.format(Name=full_name, Event=event_name, Date=event_date_str, ID=unique_id)
                    body = email_body_template.format(Name=full_name, Event=event_name, Date=event_date_str, ID=unique_id)
                    send_certificate_via_outlook(email, subject, body, individual_pdf_path)
                    status = "Sent Successfully"
                    log_lines.append(f'<span class="ts">[{ts}]</span> <span class="ok">-> SUCCESS:</span> Emailed to {email} ({unique_id})')
                else:
                    status = "Dry Run - Generated Only (Not Sent)"
                    log_lines.append(f'<span class="ts">[{ts}]</span> <span class="ok">-> DRY RUN OK:</span> PDF generated for {full_name} ({unique_id})')

            except Exception as row_error:
                status = f"ERROR: {str(row_error)}"
                log_lines.append(f'<span class="ts">[{ts}]</span> <span class="err">-> FAILED:</span> {full_name or "UNKNOWN"} — {str(row_error)}')

            results.append({
                "First Name": first_name, "Last Name": last_name, "Email ID": email,
                "Unique ID": unique_id, "Status": status,
            })

            if row_callback:
                try:
                    row_callback(status)
                except Exception:
                    pass

            flush_log()
            progress_bar.progress((idx + 1) / total)

            if idx < total - 1 and not dry_run:
                delay = random.uniform(min_delay, max_delay)
                log_lines.append(f'<span class="ts">[{datetime.now().strftime("%H:%M:%S")}]</span> <span class="warn">Throttling {delay:.1f}s (Safe Mode)...</span>')
                flush_log()
                time.sleep(delay)

        log_lines.append(f'<span class="ts">[{datetime.now().strftime("%H:%M:%S")}]</span> === Pipeline finished ===')
        flush_log()
    finally:
        try:
            st.session_state["last_run_log_lines"] = list(log_lines)  # additive: powers the Logs nav page only
        except Exception:
            pass
        shutil.rmtree(work_dir, ignore_errors=True)

    return pd.DataFrame(results)


def dataframe_to_excel_bytes(df: pd.DataFrame) -> bytes:
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Certificate Run Report")
    buffer.seek(0)
    return buffer.read()


# ==========================================================================
# NEW HELPER — single-attendee certificate preview (reuses backend as-is)
# ==========================================================================
def generate_single_certificate(row: pd.Series, idx: int, event_name: str,
                                 event_date_str: str, template_bytes: bytes) -> dict:
    """
    Generates one attendee's personalized certificate on demand, using the
    exact same personalization / QR embedding functions as the batch
    pipeline. Returns a dict with the rendered preview image bytes (if COM
    available), the PDF path, the unique ID, and the verification payload.
    """
    result = {
        "unique_id": None, "payload": None, "png_bytes": None,
        "pdf_path": None, "pptx_path": None, "error": None, "work_dir": None,
    }
    work_dir = tempfile.mkdtemp(prefix="ieee_preview_single_")
    result["work_dir"] = work_dir
    try:
        first_name = str(row.get("First Name", "")).strip()
        last_name = str(row.get("Last Name", "")).strip()
        full_name = f"{first_name} {last_name}".strip()
        unique_id = generate_unique_id(idx)
        result["unique_id"] = unique_id

        template_path = os.path.abspath(os.path.join(work_dir, "_template.pptx"))
        with open(template_path, "wb") as f:
            f.write(template_bytes)

        prs = Presentation(template_path)
        mapping = {"{{Name}}": full_name, "{{Event}}": event_name, "{{Date}}": event_date_str}
        replace_placeholders_in_pptx(prs, mapping)

        payload = build_offline_verification_payload(unique_id, full_name, event_name, event_date_str)
        result["payload"] = payload
        qr_img = generate_qr_image(payload)
        qr_path = os.path.abspath(os.path.join(work_dir, f"qr_{unique_id}.png"))
        qr_img.save(qr_path)
        insert_qr_bottom_right(prs.slides[0], qr_path, prs.slide_width, prs.slide_height)

        pptx_path = os.path.abspath(os.path.join(work_dir, f"preview_{unique_id}.pptx"))
        prs.save(pptx_path)
        result["pptx_path"] = pptx_path

        if WINDOWS_COM_AVAILABLE:
            png_path = os.path.abspath(os.path.join(work_dir, f"preview_{unique_id}.png"))
            render_slide_to_png(pptx_path, png_path, slide_index=1)
            with open(png_path, "rb") as f:
                result["png_bytes"] = f.read()

            pdf_path = os.path.abspath(os.path.join(work_dir, f"preview_{unique_id}.pdf"))
            export_pptx_to_pdf_via_powerpoint(pptx_path, pdf_path)
            result["pdf_path"] = pdf_path
    except Exception as exc:
        result["error"] = str(exc)
    return result


# ==========================================================================
# UI HELPERS
# ==========================================================================
def pill(text: str, kind: str = "neutral") -> str:
    return f'<span class="pill pill-{kind}">{text}</span>'


def compute_workflow_step() -> int:
    if st.session_state.run_completed:
        return 5
    if st.session_state.run_in_progress:
        return 5
    if st.session_state.get("_template_preview_rendered_once"):
        return 4
    if st.session_state.roster_df is not None and st.session_state.template_bytes is not None:
        return 3
    if st.session_state.roster_df is not None or st.session_state.template_bytes is not None:
        return 2
    return 1


def render_wizard(current_step: int):
    labels = ["Event Details", "Upload Files", "Validation", "Preview", "Execute"]
    html = '<div class="wizard">'
    for i, label in enumerate(labels, start=1):
        if i < current_step:
            cls, dot = "done", "✓"
        elif i == current_step:
            cls, dot = "active", str(i)
        else:
            cls, dot = "", str(i)
        html += f'<div class="wizard-step {cls}"><span class="step-dot">{dot}</span>{label}</div>'
        if i != len(labels):
            html += '<div class="wizard-arrow">›</div>'
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


def render_kpi_cards(office_ok: bool, template_status: str, roster_status: str, readiness_pct: int):
    """
    Renders all four KPI cards as ONE HTML block using a CSS grid
    (repeat(auto-fit, minmax(220px, 1fr))) so it reflows into a true
    responsive 2x2 grid on narrower widths instead of Streamlit's
    st.columns() compressing each card's text.
    """
    office_val = "Connected" if office_ok else "Unavailable"
    office_pill = pill("Ready" if office_ok else "Check Setup", "ok" if office_ok else "bad")

    template_pill_kind = "waiting" if template_status == "Not Uploaded" else ("ok" if template_status == "Valid" else "warn")
    roster_pill_kind = "waiting" if roster_status == "Not Uploaded" else "ok"

    readiness_kind = "ok" if readiness_pct == 100 else ("warn" if readiness_pct >= 50 else "bad")
    readiness_color = "var(--success)" if readiness_pct == 100 else ("var(--warning)" if readiness_pct >= 50 else "var(--error)")

    html = f"""
    <div class="kpi-grid">
        <div class="kpi-card">
            <div class="kpi-top"><span class="kpi-icon">🖥️</span>{office_pill}</div>
            <div class="kpi-label">Office COM Status</div>
            <div class="kpi-value">{office_val}</div>
            <div class="kpi-secondary">Powers PowerPoint export & Outlook send</div>
        </div>
        <div class="kpi-card">
            <div class="kpi-top"><span class="kpi-icon">🧩</span>{pill(template_status, template_pill_kind)}</div>
            <div class="kpi-label">Template Status</div>
            <div class="kpi-value">{template_status}</div>
            <div class="kpi-secondary">Checks {{Name}}, {{Event}}, {{Date}}</div>
        </div>
        <div class="kpi-card">
            <div class="kpi-top"><span class="kpi-icon">🗂️</span>{pill("Loaded" if roster_status != "Not Uploaded" else "Waiting", roster_pill_kind)}</div>
            <div class="kpi-label">Roster Status</div>
            <div class="kpi-value">{roster_status}</div>
            <div class="kpi-secondary">Requires First/Last Name + Email ID</div>
        </div>
        <div class="kpi-card">
            <div class="kpi-top"><span class="kpi-icon">📊</span>{pill(f"{readiness_pct}%", readiness_kind)}</div>
            <div class="kpi-label">Overall System Readiness</div>
            <div class="kpi-value" style="color:{readiness_color};">{readiness_pct}%</div>
            <div class="kpi-secondary">Office · Template · Roster · Placeholders</div>
        </div>
    </div>
    """
    st.markdown(html, unsafe_allow_html=True)


def render_placeholder_tiles(found_map: dict, uploaded: bool = True):
    """
    Three-state validation tiles:
      - Waiting (neutral grey): no template uploaded yet
      - Valid (green): placeholder found after upload
      - Error (red): placeholder genuinely missing after upload
    """
    cols = st.columns(len(found_map))
    for col, (key, found) in zip(cols, found_map.items()):
        with col:
            if not uploaded:
                cls, icon = "tile-waiting", "…"
            elif found:
                cls, icon = "tile-ok", "✓"
            else:
                cls, icon = "tile-bad", "✕"
            st.markdown(f'<div class="tile {cls}">{icon} {key}</div>', unsafe_allow_html=True)


@st.dialog("Certificate Preview", width="large")
def certificate_preview_dialog():
    idx = st.session_state.preview_row_idx
    roster_df = st.session_state.roster_df
    if roster_df is None or idx is None or idx not in roster_df.index:
        st.error("No attendee selected.")
        if st.button("Close", key="close_preview_missing"):
            st.session_state.show_cert_preview = False
            st.rerun()
        return

    row = roster_df.loc[idx]
    event_name = st.session_state.get("event_name_input", "").strip() or "Sample Event"
    event_date_str = st.session_state.get("event_date_str_cache", datetime.now().strftime("%d %B %Y"))

    if st.session_state.template_bytes is None:
        st.error("Upload a certificate template first.")
        if st.button("Close", key="close_preview_notemplate"):
            st.session_state.show_cert_preview = False
            st.rerun()
        return

    info_col, cert_col = st.columns([1, 2], gap="large")

    with st.spinner("Generating personalized certificate preview..."):
        preview = generate_single_certificate(row, idx, event_name, event_date_str, st.session_state.template_bytes)

    with info_col:
        full_name = f"{str(row.get('First Name','')).strip()} {str(row.get('Last Name','')).strip()}".strip()
        st.markdown("**Attendee Information**")
        st.write(f"Name: {full_name}")
        st.write(f"Email: {row.get('Email ID', '')}")
        st.write(f"Verification ID: {preview.get('unique_id')}")
        st.write(f"Event: {event_name}")
        st.write(f"Date: {event_date_str}")
        if preview.get("error"):
            st.markdown(pill("Generation Failed", "bad"), unsafe_allow_html=True)
            st.error(preview["error"])
        else:
            st.markdown(pill("QR Embedded", "ok") + " " + pill("Generated", "ok"), unsafe_allow_html=True)

        if preview.get("payload"):
            with st.expander("Offline QR Verification Payload"):
                st.code(preview["payload"], language="text")

        btn_c1, btn_c2 = st.columns(2)
        with btn_c1:
            if preview.get("pdf_path") and os.path.exists(preview["pdf_path"]):
                with open(preview["pdf_path"], "rb") as f:
                    st.download_button("Download PDF", f.read(), file_name=f"{full_name}_{preview.get('unique_id')}.pdf",
                                        mime="application/pdf", icon=":material/download:", use_container_width=True)
        with btn_c2:
            if st.button("Regenerate", icon=":material/refresh:", use_container_width=True, key="regen_preview_btn"):
                st.rerun()

        if not WINDOWS_COM_AVAILABLE:
            st.info("Live send / PDF export require Office COM automation (Windows only).")
        else:
            if st.button("Send Individual Certificate", icon=":material/mail:", use_container_width=True, key="send_single_btn"):
                try:
                    email_subject_template = st.session_state.get("email_subject_template_cache", "Your Certificate")
                    email_body_template = st.session_state.get("email_body_template_cache", "Dear {Name}, please find attached your certificate.")
                    subject = email_subject_template.format(Name=full_name, Event=event_name, Date=event_date_str, ID=preview.get("unique_id"))
                    body = email_body_template.format(Name=full_name, Event=event_name, Date=event_date_str, ID=preview.get("unique_id"))
                    send_certificate_via_outlook(str(row.get("Email ID", "")).strip(), subject, body, preview["pdf_path"])
                    st.success("Sent successfully via Outlook.")
                except Exception as send_err:
                    st.error(f"Send failed: {send_err}")

        if st.button("Close Preview", key="close_preview_btn", use_container_width=True):
            st.session_state.show_cert_preview = False
            st.rerun()

    with cert_col:
        if preview.get("png_bytes"):
            st.image(preview["png_bytes"], use_container_width=True, caption="Rendered exactly as PowerPoint would export it")
        elif preview.get("error"):
            st.warning("Preview image unavailable due to the generation error above.")
        else:
            st.info("PowerPoint COM rendering unavailable — image preview requires Windows + Office.")


# ==========================================================================
# SIDEBAR — full nav rail (Dashboard / Certificates / Monitoring / Security / Logs)
# ==========================================================================
_NAV_ITEMS = [
    ("dashboard", "🏠", "Dashboard"),
    ("certificates", "🎖️", "Certificates"),
    ("monitoring", "📈", "Monitoring"),
    ("security", "🛡️", "Security"),
    ("logs", "🖥️", "Logs"),
]

with st.sidebar:
    st.markdown(
        '<div style="padding: 4px 0 10px 0;">'
        '<div style="display:flex; align-items:center; gap:8px;">'
        '<div style="width:28px; height:28px; border-radius:8px; background:var(--accent); color:#fff; '
        'display:flex; align-items:center; justify-content:center; font-weight:800; font-size:0.85rem;">C</div>'
        '<div style="font-size:1.05rem; font-weight:800; letter-spacing:-0.01em;">Accredify</div></div>'
        '<div style="font-size:0.66rem; color:var(--text-dim); font-weight:600; text-transform:uppercase; '
        'letter-spacing:0.04em; margin-top:2px;">v2.1-Enterprise</div></div>',
        unsafe_allow_html=True,
    )
    st.markdown('<hr style="margin:2px 0 10px 0;">', unsafe_allow_html=True)

    for _nav_id, _nav_icon, _nav_label in _NAV_ITEMS:
        _is_active = st.session_state.active_nav == _nav_id
        if st.button(
            f"{_nav_icon}  {_nav_label}",
            key=f"nav_btn_{_nav_id}",
            use_container_width=True,
            type="primary" if _is_active else "secondary",
        ):
            st.session_state.active_nav = _nav_id
            st.rerun()

    st.markdown('<hr style="margin:14px 0 10px 0;">', unsafe_allow_html=True)

    _safe_mode_on = True  # Safe-Mode throttle is always enforced between sends (see Security page)
    st.markdown(
        f'<div style="padding:10px 12px; border-radius:10px; background:rgba(34,197,94,0.12); '
        f'border:1px solid rgba(34,197,94,0.30);">'
        f'<div style="font-size:0.78rem; font-weight:700; color:#4ADE80;">🔒 Safe-Mode Active</div>'
        f'<div style="font-size:0.68rem; color:var(--text-dim); margin-top:2px;">Throttled sends between attendees</div>'
        f'</div>',
        unsafe_allow_html=True,
    )

    st.markdown(
        pill("PowerPoint COM: Ready", "ok") if WINDOWS_COM_AVAILABLE else pill("PowerPoint COM: Down", "bad"),
        unsafe_allow_html=True,
    )
    st.markdown(
        pill("Outlook COM: Ready", "ok") if WINDOWS_COM_AVAILABLE else pill("Outlook COM: Down", "bad"),
        unsafe_allow_html=True,
    )
    st.markdown(
        pill(f"Env: {'Windows' if IS_WINDOWS else 'Non-Windows'}", "neutral"),
        unsafe_allow_html=True,
    )

    st.markdown('<hr style="margin:12px 0;">', unsafe_allow_html=True)
    if st.button("⚙️  Settings", key="nav_btn_settings", use_container_width=True, type="secondary"):
        st.session_state.active_nav = "security"
        st.rerun()
    st.caption("❓ Support · IEEE Student Branch VIT")
    st.caption("Local-only automation · No external mail relay")


# ==========================================================================
# TOP NAVIGATION HEADER
# ==========================================================================
outlook_ready = WINDOWS_COM_AVAILABLE  # same COM bridge powers both apps locally
_NAV_TITLES = {
    "dashboard": "Accredify Certificate Automation Suite",
    "certificates": "Generated Certificates",
    "monitoring": "Pipeline Monitoring",
    "security": "Security & Configuration",
    "logs": "System Journal",
}
st.markdown(
    f"""
    <div class="topnav">
        <div class="topnav-brand">🎓 Accredify</div>
        <div class="topnav-title">{_NAV_TITLES.get(st.session_state.active_nav, "Accredify Certificate Automation Suite")}</div>
        <div class="topnav-status">
            {pill("Office COM Connected", "ok") if WINDOWS_COM_AVAILABLE else pill("Office COM Unavailable", "bad")}
            {pill("Outlook Ready", "ok") if outlook_ready else pill("Outlook Not Ready", "warn")}
            {pill("Offline Verification", "neutral")}
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

if st.session_state.active_nav == "dashboard":
    render_wizard(compute_workflow_step())


if st.session_state.active_nav == "dashboard":
    # ==========================================================================
    # HELP & DOCUMENTATION (moved out of sidebar into a collapsible accordion)
    # ==========================================================================
    with st.expander("📘 Help & Documentation", expanded=False):
        st.markdown(
            """
            **Workflow:** Event Details → Upload Files → Validation → Preview → Execute

            - Roster must be `.xlsx` with columns exactly: `First Name`, `Last Name`, `Email ID`
            - Template must be `.pptx` with placeholders `{{Name}}`, `{{Event}}`, `{{Date}}` on slide 1
            - The bottom-right 1.3"×1.3" corner (0.3" margin) is reserved for the QR code — keep it clear
            - Requires Windows + classic Outlook (not the new standalone Outlook app) + PowerPoint, both already licensed and signed in
            - Use **Dry Run Mode** to generate PDFs locally without sending any email — recommended before your first live batch
            - The offline QR payload contains the full verification card text directly — no server, no link, works with any camera app
            """
        )

    tab_dashboard, tab_layout = st.tabs(["📊 Execution Dashboard", "🖼️ Layout Preview Workspace"])


    # ==========================================================================
    # TAB 1 — EXECUTION DASHBOARD
    # ==========================================================================
    with tab_dashboard:
        min_delay = st.session_state.get("min_delay_slider", 5)
        max_delay = st.session_state.get("max_delay_slider", 12)
        if max_delay < min_delay:
            max_delay = min_delay

        left_col, right_col = st.columns([7, 3], gap="large")

        # ---------------- LEFT COLUMN ----------------
        with left_col:
            st.markdown('<div class="card-primary"><div class="card-title">📌 Event Details</div>', unsafe_allow_html=True)
            ev_c1, ev_c2 = st.columns(2)
            with ev_c1:
                event_name_input = st.text_input("Event Name", placeholder="e.g., Accredify TechFest 2026", key="event_name_input")
            with ev_c2:
                event_date_input = st.date_input("Event Date", value=date.today(), key="event_date_input")
            event_date_str = event_date_input.strftime("%d %B %Y")
            st.session_state["event_date_str_cache"] = event_date_str
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown('<div class="card-primary"><div class="card-title">📂 Upload Files</div>', unsafe_allow_html=True)
            up_c1, up_c2 = st.columns(2)
            with up_c1:
                st.markdown(
                    '<div class="upload-illustration"><div class="upicon">📄⬆️</div>'
                    '<b>Attendee Roster</b><div class="upload-meta">Drag & drop or browse · .xlsx · max 10MB</div></div>',
                    unsafe_allow_html=True,
                )
                roster_file = st.file_uploader("Roster", type=["xlsx"], label_visibility="collapsed", key="roster_uploader")
            with up_c2:
                st.markdown(
                    '<div class="upload-illustration"><div class="upicon">🖼️⬆️</div>'
                    '<b>Certificate Template</b><div class="upload-meta">Drag & drop or browse · .pptx · max 25MB</div></div>',
                    unsafe_allow_html=True,
                )
                pptx_template_file = st.file_uploader("Template", type=["pptx"], label_visibility="collapsed", key="template_uploader")
            st.markdown("</div>", unsafe_allow_html=True)

            roster_df = None
            if roster_file is not None:
                try:
                    roster_df = pd.read_excel(roster_file, engine="openpyxl")
                    required_cols = {"First Name", "Last Name", "Email ID"}
                    missing_cols = required_cols - set(roster_df.columns)
                    if missing_cols:
                        st.error(f"Missing required column(s): {', '.join(sorted(missing_cols))}")
                        roster_df = None
                    else:
                        roster_df = roster_df.dropna(how="all").reset_index(drop=True)
                        st.session_state.roster_df = roster_df
                        st.markdown(
                            f'<div class="upload-success">✓ {roster_file.name} · {len(roster_df)} attendee(s) validated</div>',
                            unsafe_allow_html=True,
                        )
                except Exception as parse_error:
                    st.error(f"Failed to parse Excel roster: {parse_error}")
                    roster_df = None
            elif st.session_state.roster_df is not None:
                roster_df = st.session_state.roster_df
                st.caption(f"Using previously uploaded roster · {len(roster_df)} attendee(s)")

            template_bytes = None
            found_map = {"{{Name}}": False, "{{Event}}": False, "{{Date}}": False}
            if pptx_template_file is not None:
                template_bytes = pptx_template_file.getvalue()
                st.session_state.template_bytes = template_bytes
                st.session_state.template_name = pptx_template_file.name
                found_map = pptx_has_all_placeholders(template_bytes)
                if all(found_map.values()):
                    st.markdown(
                        f'<div class="upload-success">✓ {pptx_template_file.name} · all placeholders detected</div>',
                        unsafe_allow_html=True,
                    )
                else:
                    missing = [k for k, v in found_map.items() if not v]
                    st.warning(f"Template uploaded but missing: {', '.join(missing)}")
            elif st.session_state.template_bytes is not None:
                template_bytes = st.session_state.template_bytes
                found_map = pptx_has_all_placeholders(template_bytes)
                st.caption(f"Using previously uploaded template · {st.session_state.template_name}")

            with st.expander("✉️ Advanced: Email Subject & Body Templates", expanded=False):
                email_subject_template = st.text_input(
                    "Email Subject",
                    value=st.session_state.get("email_subject_template_cache", "Your Certificate for {Event} — IEEE Student Branch VIT"),
                    key="email_subject_template_input",
                )
                email_body_template = st.text_area(
                    "Email Body", height=160,
                    value=st.session_state.get("email_body_template_cache", (
                        "Dear {Name},\n\nThank you for participating in {Event} held on {Date}, organized by "
                        "the IEEE Student Branch, VIT.\n\nYour certificate is attached as a PDF and carries a "
                        "unique offline-verifiable QR code (ID: {ID}) — scan with any camera app, no internet "
                        "required.\n\nBest regards,\nIEEE Student Branch, VIT"
                    )),
                    key="email_body_template_input",
                )
                st.session_state["email_subject_template_cache"] = email_subject_template
                st.session_state["email_body_template_cache"] = email_body_template

            dry_run_toggle = st.session_state.get("dry_run_toggle", False)
            st.caption(
                f"🧪 Dry Run Mode is **{'ON' if dry_run_toggle else 'OFF'}** — change this on the "
                f"**Security** page in the left nav."
            )

            # ---------------- Attendee management table ----------------
            st.markdown('<div class="card-primary"><div class="card-title">🗂️ Attendee Management</div>', unsafe_allow_html=True)
            if roster_df is None or len(roster_df) == 0:
                st.info("Upload a roster to see the interactive attendee table here.")
            else:
                tc1, tc2, tc3 = st.columns([2, 1, 1])
                with tc1:
                    search_query = st.text_input("Search by name or email", value=st.session_state.attendee_search, key="attendee_search_input")
                with tc2:
                    sort_choice = st.selectbox("Sort by", ["Row order", "First Name", "Last Name", "Email ID"], key="attendee_sort_input")
                with tc3:
                    status_filter = st.selectbox("Status filter", ["All", "Valid email", "Invalid email"], key="attendee_status_filter_input")

                display_df = roster_df.copy()
                display_df["__idx"] = display_df.index
                if search_query:
                    mask = (
                        display_df["First Name"].astype(str).str.contains(search_query, case=False, na=False)
                        | display_df["Last Name"].astype(str).str.contains(search_query, case=False, na=False)
                        | display_df["Email ID"].astype(str).str.contains(search_query, case=False, na=False)
                    )
                    display_df = display_df[mask]
                valid_email_mask = display_df["Email ID"].astype(str).str.contains("@", na=False)
                if status_filter == "Valid email":
                    display_df = display_df[valid_email_mask]
                elif status_filter == "Invalid email":
                    display_df = display_df[~valid_email_mask]
                if sort_choice != "Row order":
                    display_df = display_df.sort_values(by=sort_choice, na_position="last")

                # ---- Pagination ----
                PAGE_SIZE = 8
                total_filtered = len(display_df)
                total_pages = max(1, (total_filtered - 1) // PAGE_SIZE + 1)
                if "attendee_page" not in st.session_state:
                    st.session_state.attendee_page = 1
                st.session_state.attendee_page = min(st.session_state.attendee_page, total_pages)
                start_i = (st.session_state.attendee_page - 1) * PAGE_SIZE
                page_df = display_df.iloc[start_i:start_i + PAGE_SIZE]

                st.caption(f"{total_filtered} of {len(roster_df)} attendee(s) shown · Columns detected: {', '.join(roster_df.columns)}")

                st.markdown(
                    '<div class="tbl-header"><div>Name</div><div>Email</div><div>Status</div><div>Actions</div></div>',
                    unsafe_allow_html=True,
                )
                for _, r in page_df.iterrows():
                    idx = int(r["__idx"])
                    full_name = f"{str(r.get('First Name','')).strip()} {str(r.get('Last Name','')).strip()}".strip()
                    email_val = str(r.get("Email ID", "")).strip()
                    is_valid = "@" in email_val and "." in email_val.split("@")[-1] if email_val else False
                    row_c1, row_c2, row_c3, row_c4 = st.columns([2, 2.2, 1.1, 2])
                    row_c1.write(full_name or "—")
                    row_c2.write(email_val or "—")
                    row_c3.markdown(pill("Valid", "ok") if is_valid else pill("Invalid", "bad"), unsafe_allow_html=True)
                    with row_c4:
                        act1, act2, act3, act4 = st.columns(4)
                        with act1:
                            if st.button("", key=f"preview_btn_{idx}", icon=":material/visibility:", help="Preview Certificate", use_container_width=True):
                                st.session_state.preview_row_idx = idx
                                st.session_state.show_cert_preview = True
                                st.rerun()
                        with act2:
                            dl_disabled = template_bytes is None or not WINDOWS_COM_AVAILABLE
                            if st.button("", key=f"dl_btn_{idx}", icon=":material/download:", help="Download Certificate", use_container_width=True, disabled=dl_disabled):
                                preview = generate_single_certificate(r, idx, event_name_input.strip() or "Sample Event", event_date_str, template_bytes)
                                if preview.get("error"):
                                    st.error(preview["error"])
                                elif preview.get("pdf_path") and os.path.exists(preview["pdf_path"]):
                                    st.session_state[f"_dl_ready_{idx}"] = open(preview["pdf_path"], "rb").read()
                                    st.session_state[f"_dl_name_{idx}"] = f"{full_name}_{preview['unique_id']}.pdf"
                        with act3:
                            if st.button("", key=f"send_btn_{idx}", icon=":material/mail:", help="Send Individual Certificate", use_container_width=True,
                                         disabled=not is_valid or template_bytes is None or not WINDOWS_COM_AVAILABLE):
                                preview = generate_single_certificate(r, idx, event_name_input.strip() or "Sample Event", event_date_str, template_bytes)
                                if preview.get("error"):
                                    st.error(preview["error"])
                                else:
                                    try:
                                        subj = email_subject_template.format(Name=full_name, Event=event_name_input, Date=event_date_str, ID=preview["unique_id"])
                                        bod = email_body_template.format(Name=full_name, Event=event_name_input, Date=event_date_str, ID=preview["unique_id"])
                                        send_certificate_via_outlook(email_val, subj, bod, preview["pdf_path"])
                                        st.success(f"Sent to {email_val}")
                                    except Exception as e:
                                        st.error(str(e))
                        with act4:
                            if st.button("", key=f"regen_btn_{idx}", icon=":material/refresh:", help="Regenerate Certificate", use_container_width=True,
                                         disabled=template_bytes is None):
                                st.session_state.preview_row_idx = idx
                                st.session_state.show_cert_preview = True
                                st.rerun()
                    if st.session_state.get(f"_dl_ready_{idx}"):
                        st.download_button(
                            f"Download PDF — {full_name}", data=st.session_state[f"_dl_ready_{idx}"],
                            file_name=st.session_state[f"_dl_name_{idx}"], mime="application/pdf",
                            key=f"dl_confirm_{idx}", use_container_width=True,
                        )

                pg_c1, pg_c2, pg_c3 = st.columns([1, 2, 1])
                with pg_c1:
                    if st.button("‹ Prev", key="attendee_prev_page", disabled=st.session_state.attendee_page <= 1, type="secondary"):
                        st.session_state.attendee_page -= 1
                        st.rerun()
                with pg_c2:
                    st.markdown(f'<div class="tbl-pagination" style="justify-content:center;">Page {st.session_state.attendee_page} of {total_pages}</div>', unsafe_allow_html=True)
                with pg_c3:
                    if st.button("Next ›", key="attendee_next_page", disabled=st.session_state.attendee_page >= total_pages, type="secondary"):
                        st.session_state.attendee_page += 1
                        st.rerun()

            if st.session_state.show_cert_preview:
                certificate_preview_dialog()

            st.markdown("</div>", unsafe_allow_html=True)

            # ---------------- Execution console + dashboard ----------------
            st.markdown('<div class="card"><div class="card-title">🖥️ Execution Console</div>', unsafe_allow_html=True)
            stage_placeholder = st.empty()
            eta_placeholder = st.empty()
            progress_bar_placeholder = st.empty()
            console_log_placeholder = st.empty()
            console_log_placeholder.markdown('<div class="terminal">Awaiting execution...</div>', unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)

            # ---------------- Sticky-ish bottom action bar ----------------
            st.markdown('<div class="action-bar">', unsafe_allow_html=True)
            bar_c1, bar_c2, bar_c3, bar_c4 = st.columns([1.3, 1.3, 1, 1.6])
            with bar_c1:
                st.write("")  # dry run toggle already placed above; this bar focuses on actions
                st.caption(f"Dry Run: {'ON' if dry_run_toggle else 'OFF'}")
            with bar_c2:
                preview_template_clicked = st.button("Preview Template", icon=":material/slideshow:", use_container_width=True)
            with bar_c3:
                reset_clicked = st.button("Reset Form", icon=":material/restart_alt:", use_container_width=True, type="secondary")
            with bar_c4:
                run_clicked = st.button("🚀 Execute Certificate Pipeline", use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

            if reset_clicked:
                for key in ["run_results_df", "run_completed", "roster_df", "template_bytes", "template_name"]:
                    st.session_state[key] = _defaults[key]
                st.rerun()

            if preview_template_clicked:
                if template_bytes is None:
                    st.error("Upload a certificate template first.")
                elif not WINDOWS_COM_AVAILABLE:
                    st.error("Template rendering requires PowerPoint COM automation on Windows.")
                else:
                    preview_work_dir = tempfile.mkdtemp(prefix="ieee_tpl_preview_")
                    try:
                        with st.spinner("Rendering slide 1 via PowerPoint..."):
                            tpl_path = os.path.abspath(os.path.join(preview_work_dir, "tpl.pptx"))
                            with open(tpl_path, "wb") as f:
                                f.write(template_bytes)
                            png_path = os.path.abspath(os.path.join(preview_work_dir, "tpl.png"))
                            render_slide_to_png(tpl_path, png_path, slide_index=1)
                            with open(png_path, "rb") as f:
                                st.session_state["_quick_template_png"] = f.read()
                        st.session_state["_template_preview_rendered_once"] = True
                        st.image(st.session_state["_quick_template_png"], caption="Slide 1 preview", use_container_width=True)
                    except Exception as e:
                        st.error(f"Preview failed: {e}")
                    finally:
                        shutil.rmtree(preview_work_dir, ignore_errors=True)

            if run_clicked:
                validation_errors = []
                if not IS_WINDOWS or not WINDOWS_COM_AVAILABLE:
                    validation_errors.append("Office COM automation is unavailable. Run on Windows with PowerPoint and classic Outlook installed.")
                if not event_name_input.strip():
                    validation_errors.append("Event Name cannot be empty.")
                if roster_df is None or len(roster_df) == 0:
                    validation_errors.append("A valid attendee roster (.xlsx) must be uploaded.")
                if template_bytes is None:
                    validation_errors.append("A certificate template (.pptx) must be uploaded.")

                if validation_errors:
                    for err in validation_errors:
                        st.error(err)
                else:
                    st.session_state.run_in_progress = True
                    progress_bar = progress_bar_placeholder.progress(0.0)
                    run_start = time.time()
                    counters = {"processed": 0, "success": 0, "failed": 0}
                    total_rows = len(roster_df)

                    def _stage_cb(stage_name, idx, total, full_name):
                        stage_placeholder.markdown(
                            f"**Current Stage:** `{stage_name}`  ·  **Attendee:** {full_name} ({idx + 1}/{total})"
                        )

                    def _row_cb(status):
                        counters["processed"] += 1
                        if status.startswith("Sent Successfully") or status.startswith("Dry Run"):
                            counters["success"] += 1
                        else:
                            counters["failed"] += 1
                        elapsed = time.time() - run_start
                        avg = elapsed / max(counters["processed"], 1)
                        remaining = total_rows - counters["processed"]
                        eta_seconds = max(int(avg * remaining), 0)
                        eta_placeholder.markdown(
                            f"Processed: **{counters['processed']}/{total_rows}**  ·  "
                            f"✅ Success: **{counters['success']}**  ·  "
                            f"❌ Failed: **{counters['failed']}**  ·  "
                            f"⏱ Est. remaining: **~{eta_seconds}s**"
                        )

                    with st.spinner("Running certificate pipeline — do not close this tab..."):
                        result_df = process_roster(
                            df=roster_df,
                            event_name=event_name_input.strip(),
                            event_date_str=event_date_str,
                            pptx_template_bytes=template_bytes,
                            min_delay=float(min_delay),
                            max_delay=float(max_delay),
                            log_placeholder=console_log_placeholder,
                            progress_bar=progress_bar,
                            email_subject_template=email_subject_template,
                            email_body_template=email_body_template,
                            dry_run=dry_run_toggle,
                            stage_callback=_stage_cb,
                            row_callback=_row_cb,
                        )
                    st.session_state.run_results_df = result_df
                    st.session_state.run_completed = True
                    st.session_state.run_in_progress = False

            # ---------------- Run summary analytics ----------------
            if st.session_state.run_completed and st.session_state.run_results_df is not None:
                result_df = st.session_state.run_results_df
                success_count = int(result_df["Status"].astype(str).str.startswith("Sent Successfully").sum())
                dry_run_count = int(result_df["Status"].astype(str).str.startswith("Dry Run").sum())
                error_count = int(result_df["Status"].astype(str).str.startswith("ERROR").sum())
                total_count = len(result_df)
                success_rate = round(((success_count + dry_run_count) / total_count) * 100, 1) if total_count else 0.0

                st.markdown('<div class="card"><div class="card-title">📊 Batch Completion Summary</div>', unsafe_allow_html=True)
                m1, m2, m3, m4 = st.columns(4)
                m1.metric("Total Attendees", total_count)
                m2.metric("Certificates Generated", success_count + dry_run_count)
                m3.metric("Emails Sent", success_count)
                m4.metric("Failed Records", error_count)
                m5, m6 = st.columns(2)
                m5.metric("Success Rate", f"{success_rate}%")
                m6.metric("Dry Run Generated", dry_run_count)

                st.dataframe(result_df, use_container_width=True, height=280)
                excel_bytes = dataframe_to_excel_bytes(result_df)
                st.download_button(
                    "⬇️ Download Full Status Report (.xlsx)", data=excel_bytes,
                    file_name=f"IEEE_VIT_Certificate_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    icon=":material/download:", use_container_width=True,
                )
                st.markdown("</div>", unsafe_allow_html=True)

        # ---------------- RIGHT COLUMN — monitoring cards ----------------
        with right_col:
            template_status_label = "Not Uploaded"
            if template_bytes is not None:
                template_status_label = "Valid" if all(found_map.values()) else "Needs Review"
            roster_status_label = "Not Uploaded"
            if roster_df is not None:
                roster_status_label = f"{len(roster_df)} Loaded"

            checks = [WINDOWS_COM_AVAILABLE, template_bytes is not None, roster_df is not None, all(found_map.values()) if template_bytes else False]
            readiness_pct = int(round(sum(1 for c in checks if c) / len(checks) * 100))

            render_kpi_cards(WINDOWS_COM_AVAILABLE, template_status_label, roster_status_label, readiness_pct)

            st.markdown('<div class="card"><div class="card-title">🧭 System Status</div>', unsafe_allow_html=True)
            st.markdown(pill("PowerPoint COM", "ok" if WINDOWS_COM_AVAILABLE else "bad"), unsafe_allow_html=True)
            st.markdown(pill("Outlook COM", "ok" if WINDOWS_COM_AVAILABLE else "bad"), unsafe_allow_html=True)
            st.markdown(pill("Offline QR Engine", "ok"), unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown('<div class="card"><div class="card-title">✅ Validation Overview</div>', unsafe_allow_html=True)
            st.write(f"Roster: {roster_status_label}")
            st.write(f"Template: {template_status_label}")
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown('<div class="card"><div class="card-title">🧩 Placeholder Validation</div>', unsafe_allow_html=True)
            render_placeholder_tiles(found_map)
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown('<div class="card"><div class="card-title">📈 Quick Metrics</div>', unsafe_allow_html=True)
            if st.session_state.run_results_df is not None:
                rdf = st.session_state.run_results_df
                st.write(f"Last run: {len(rdf)} attendee(s)")
                st.write(f"Success: {int(rdf['Status'].astype(str).str.startswith('Sent Successfully').sum())}")
                st.write(f"Failed: {int(rdf['Status'].astype(str).str.startswith('ERROR').sum())}")
            else:
                st.caption("No run yet — metrics will appear here live during execution.")
            st.markdown("</div>", unsafe_allow_html=True)


    # ==========================================================================
    # TAB 2 — LAYOUT PREVIEW WORKSPACE
    # ==========================================================================
    with tab_layout:
        preview_main, preview_side = st.columns([7, 3], gap="large")

        with preview_main:
            st.markdown('<div class="card"><div class="card-title">🖼️ Certificate Layout Preview</div>', unsafe_allow_html=True)
            lp_upload_col, lp_ctrl_col = st.columns([2, 1])
            with lp_upload_col:
                layout_pptx_file = st.file_uploader("Certificate Template for Preview (.pptx)", type=["pptx"], key="layout_preview_uploader")
            with lp_ctrl_col:
                st.session_state.template_preview_zoom = st.slider("Zoom", 50, 150, st.session_state.template_preview_zoom, step=10, key="zoom_slider")

            ctrl_c1, ctrl_c2, ctrl_c3 = st.columns(3)
            with ctrl_c1:
                render_clicked = st.button("Render Preview", icon=":material/play_arrow:", use_container_width=True)
            with ctrl_c2:
                zoom_in_clicked = st.button("Zoom In", icon=":material/zoom_in:", use_container_width=True)
            with ctrl_c3:
                zoom_out_clicked = st.button("Zoom Out", icon=":material/zoom_out:", use_container_width=True)

            if zoom_in_clicked:
                st.session_state.template_preview_zoom = min(150, st.session_state.template_preview_zoom + 10)
            if zoom_out_clicked:
                st.session_state.template_preview_zoom = max(50, st.session_state.template_preview_zoom - 10)

            if render_clicked:
                if layout_pptx_file is None:
                    st.error("Please upload a .pptx template first.")
                elif not WINDOWS_COM_AVAILABLE:
                    st.error("Live rendering requires PowerPoint COM automation on Windows — unavailable on this system.")
                else:
                    preview_work_dir = tempfile.mkdtemp(prefix="ieee_layout_preview_")
                    try:
                        with st.spinner("Rendering slide 1 via PowerPoint..."):
                            pptx_preview_path = os.path.abspath(os.path.join(preview_work_dir, "preview_template.pptx"))
                            with open(pptx_preview_path, "wb") as f:
                                f.write(layout_pptx_file.getvalue())
                            png_output_path = os.path.abspath(os.path.join(preview_work_dir, "preview_slide.png"))
                            render_slide_to_png(pptx_preview_path, png_output_path, slide_index=1)
                            with open(png_output_path, "rb") as f:
                                st.session_state["_layout_preview_png"] = f.read()
                    except Exception as preview_error:
                        st.error(f"Preview rendering failed: {preview_error}")
                    finally:
                        shutil.rmtree(preview_work_dir, ignore_errors=True)

            if st.session_state.get("_layout_preview_png"):
                zoom_pct = st.session_state.template_preview_zoom
                st.image(st.session_state["_layout_preview_png"], width=int(1000 * zoom_pct / 100), caption=f"Slide 1 — {zoom_pct}% zoom")
            else:
                st.info("Upload a template and click Render Preview to see a real PowerPoint-rendered image here.")
            st.markdown("</div>", unsafe_allow_html=True)

        with preview_side:
            st.markdown('<div class="card"><div class="card-title">🔳 QR Preview</div>', unsafe_allow_html=True)
            sample_event = st.session_state.get("event_name_input", "").strip() or "Sample Event Name"
            sample_date = st.session_state.get("event_date_str_cache", "01 January 2026")
            sample_payload = build_offline_verification_payload("IEEE-VIT-2026-0001", "Jordan Example", sample_event, sample_date)
            st.code(sample_payload, language="text")
            sample_qr_img = generate_qr_image(sample_payload)
            qr_buffer = io.BytesIO()
            sample_qr_img.save(qr_buffer, format="PNG")
            st.image(qr_buffer.getvalue(), width=160)
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown('<div class="card"><div class="card-title">🧩 Placeholder Detection</div>', unsafe_allow_html=True)
            if layout_pptx_file is not None:
                layout_found_map = pptx_has_all_placeholders(layout_pptx_file.getvalue())
                render_placeholder_tiles(layout_found_map)
            else:
                st.caption("Upload a template above to check placeholders.")
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown('<div class="card"><div class="card-title">✅ Pre-Flight Checklist</div>', unsafe_allow_html=True)
            st.checkbox("Placeholders {{Name}}, {{Event}}, {{Date}} present & correctly spelled", key="chk_placeholders")
            st.checkbox("Bottom-right corner clear (1.3\" × 1.3\", 0.3\" margin)", key="chk_corner")
            st.checkbox("Roster columns exactly: First Name, Last Name, Email ID", key="chk_columns")
            st.checkbox("Classic Outlook is open and signed in", key="chk_outlook")
            st.checkbox("Safe-Mode throttle set appropriately for roster size", key="chk_throttle")
            st.markdown("</div>", unsafe_allow_html=True)

elif st.session_state.active_nav == "certificates":
    st.markdown('<div class="card-primary"><div class="card-title">🎖️ Generated Certificates</div>'
                '<div class="card-subtle">Reuses the exact same generation, QR, and Outlook-send backend as the Dashboard — nothing is re-implemented here.</div>',
                unsafe_allow_html=True)
    _cert_roster = st.session_state.roster_df
    _cert_template = st.session_state.template_bytes
    if _cert_roster is None or len(_cert_roster) == 0:
        st.info("No roster loaded yet. Upload one from the **Dashboard** page to see certificates here.")
    elif _cert_template is None:
        st.info("No certificate template loaded yet. Upload one from the **Dashboard** page first.")
    else:
        _cert_event_name = st.session_state.get("event_name_input", "").strip() or "Sample Event"
        _cert_event_date = st.session_state.get("event_date_str_cache", datetime.now().strftime("%d %B %Y"))
        st.markdown(
            '<div class="tbl-header"><div>Name</div><div>Email</div><div>Status</div><div>Actions</div></div>',
            unsafe_allow_html=True,
        )
        for _c_idx, _c_row in _cert_roster.iterrows():
            _c_name = f"{str(_c_row.get('First Name','')).strip()} {str(_c_row.get('Last Name','')).strip()}".strip()
            _c_email = str(_c_row.get("Email ID", "")).strip()
            _c_valid = "@" in _c_email and "." in _c_email.split("@")[-1] if _c_email else False
            _cc1, _cc2, _cc3, _cc4 = st.columns([2, 2.2, 1.1, 2])
            _cc1.write(_c_name or "—")
            _cc2.write(_c_email or "—")
            _cc3.markdown(pill("Valid", "ok") if _c_valid else pill("Invalid", "bad"), unsafe_allow_html=True)
            with _cc4:
                _a1, _a2, _a3 = st.columns(3)
                with _a1:
                    if st.button("", key=f"cert_preview_{_c_idx}", icon=":material/visibility:", help="Preview Certificate", use_container_width=True):
                        st.session_state.preview_row_idx = _c_idx
                        st.session_state.show_cert_preview = True
                        st.rerun()
                with _a2:
                    if st.button("", key=f"cert_download_{_c_idx}", icon=":material/download:", help="Download Certificate",
                                 use_container_width=True, disabled=not WINDOWS_COM_AVAILABLE):
                        _preview = generate_single_certificate(_c_row, _c_idx, _cert_event_name, _cert_event_date, _cert_template)
                        if _preview.get("error"):
                            st.error(_preview["error"])
                        elif _preview.get("pdf_path") and os.path.exists(_preview["pdf_path"]):
                            st.session_state[f"_cert_dl_ready_{_c_idx}"] = open(_preview["pdf_path"], "rb").read()
                            st.session_state[f"_cert_dl_name_{_c_idx}"] = f"{_c_name}_{_preview['unique_id']}.pdf"
                with _a3:
                    if st.button("", key=f"cert_send_{_c_idx}", icon=":material/mail:", help="Resend Certificate",
                                 use_container_width=True, disabled=not _c_valid or not WINDOWS_COM_AVAILABLE):
                        _preview = generate_single_certificate(_c_row, _c_idx, _cert_event_name, _cert_event_date, _cert_template)
                        if _preview.get("error"):
                            st.error(_preview["error"])
                        else:
                            try:
                                _subj_tpl = st.session_state.get("email_subject_template_cache", "Your Certificate")
                                _body_tpl = st.session_state.get("email_body_template_cache", "Dear {Name}, please find attached your certificate.")
                                _subj = _subj_tpl.format(Name=_c_name, Event=_cert_event_name, Date=_cert_event_date, ID=_preview.get("unique_id"))
                                _body = _body_tpl.format(Name=_c_name, Event=_cert_event_name, Date=_cert_event_date, ID=_preview.get("unique_id"))
                                send_certificate_via_outlook(_c_email, _subj, _body, _preview["pdf_path"])
                                st.success(f"Sent to {_c_email}")
                            except Exception as _send_err:
                                st.error(str(_send_err))
            if st.session_state.get(f"_cert_dl_ready_{_c_idx}"):
                st.download_button(
                    f"Download PDF — {_c_name}", data=st.session_state[f"_cert_dl_ready_{_c_idx}"],
                    file_name=st.session_state[f"_cert_dl_name_{_c_idx}"], mime="application/pdf",
                    key=f"cert_dl_confirm_{_c_idx}", use_container_width=True,
                )
    if st.session_state.show_cert_preview:
        certificate_preview_dialog()
    st.markdown("</div>", unsafe_allow_html=True)

    if st.session_state.run_results_df is not None:
        st.markdown('<div class="card"><div class="card-title">📄 Last Run — Status Report</div>', unsafe_allow_html=True)
        st.dataframe(st.session_state.run_results_df, use_container_width=True, height=280)
        st.markdown("</div>", unsafe_allow_html=True)


elif st.session_state.active_nav == "monitoring":
    _mon_roster = st.session_state.roster_df
    _mon_template = st.session_state.template_bytes
    _mon_template_status = "Not Uploaded"
    _mon_found_map = {"{{Name}}": False, "{{Event}}": False, "{{Date}}": False}
    if _mon_template is not None:
        _mon_found_map = pptx_has_all_placeholders(_mon_template)
        _mon_template_status = "Valid" if all(_mon_found_map.values()) else "Needs Review"
    _mon_roster_status = "Not Uploaded" if _mon_roster is None else f"{len(_mon_roster)} Loaded"
    _mon_checks = [WINDOWS_COM_AVAILABLE, _mon_template is not None, _mon_roster is not None,
                   all(_mon_found_map.values()) if _mon_template else False]
    _mon_readiness_pct = int(round(sum(1 for c in _mon_checks if c) / len(_mon_checks) * 100))

    render_kpi_cards(WINDOWS_COM_AVAILABLE, _mon_template_status, _mon_roster_status, _mon_readiness_pct)

    _mon_c1, _mon_c2 = st.columns(2, gap="large")
    with _mon_c1:
        st.markdown('<div class="card"><div class="card-title">🧭 Live System Status</div>', unsafe_allow_html=True)
        st.markdown(pill("PowerPoint Automation", "ok" if WINDOWS_COM_AVAILABLE else "bad"), unsafe_allow_html=True)
        st.markdown(pill("Outlook COM Bridge", "ok" if WINDOWS_COM_AVAILABLE else "bad"), unsafe_allow_html=True)
        st.markdown(pill("Offline QR Engine", "ok"), unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="card"><div class="card-title">🧩 Placeholder Validation</div>', unsafe_allow_html=True)
        render_placeholder_tiles(_mon_found_map, uploaded=_mon_template is not None)
        st.markdown("</div>", unsafe_allow_html=True)

    with _mon_c2:
        st.markdown('<div class="card"><div class="card-title">📈 Quick Metrics</div>', unsafe_allow_html=True)
        if st.session_state.run_results_df is not None:
            _rdf = st.session_state.run_results_df
            st.write(f"Last run: {len(_rdf)} attendee(s)")
            st.write(f"Success: {int(_rdf['Status'].astype(str).str.startswith('Sent Successfully').sum())}")
            st.write(f"Dry Run Generated: {int(_rdf['Status'].astype(str).str.startswith('Dry Run').sum())}")
            st.write(f"Failed: {int(_rdf['Status'].astype(str).str.startswith('ERROR').sum())}")
        else:
            st.caption("No run yet — metrics will appear here after you execute the pipeline from the Dashboard.")
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="card"><div class="card-title">✅ Validation Overview</div>', unsafe_allow_html=True)
        st.write(f"Roster: {_mon_roster_status}")
        st.write(f"Template: {_mon_template_status}")
        st.markdown("</div>", unsafe_allow_html=True)


elif st.session_state.active_nav == "security":
    st.markdown('<div class="card-primary"><div class="card-title">🛡️ Security & Configuration</div>'
                '<div class="card-subtle">These controls feed directly into the same process_roster() pipeline used on the Dashboard — moving them here does not change how they behave.</div>',
                unsafe_allow_html=True)

    st.markdown('<div style="font-size:0.8rem; font-weight:700; color:var(--text-dim); text-transform:uppercase; margin-bottom:6px;">Safe-Mode Throttle</div>', unsafe_allow_html=True)
    st.slider("Min delay between sends (s)", 1, 30, 5, key="min_delay_slider")
    st.slider("Max delay between sends (s)", 1, 60, 12, key="max_delay_slider")
    st.caption("Randomized delay applied between each live send to avoid triggering Outlook/anti-spam throttling.")

    st.markdown('<hr style="margin:16px 0;">', unsafe_allow_html=True)
    st.markdown('<div style="font-size:0.8rem; font-weight:700; color:var(--text-dim); text-transform:uppercase; margin-bottom:6px;">Execution Mode</div>', unsafe_allow_html=True)
    st.checkbox(
        "🧪 Dry Run Mode (generate PDFs locally, do NOT send emails)",
        value=st.session_state.get("dry_run_toggle", False), key="dry_run_toggle",
    )
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown('<div class="card"><div class="card-title">🧭 Office COM Status</div>', unsafe_allow_html=True)
    st.markdown(pill("PowerPoint COM: Ready", "ok") if WINDOWS_COM_AVAILABLE else pill("PowerPoint COM: Down", "bad"), unsafe_allow_html=True)
    st.markdown(pill("Outlook COM: Ready", "ok") if WINDOWS_COM_AVAILABLE else pill("Outlook COM: Down", "bad"), unsafe_allow_html=True)
    st.markdown(pill(f"Env: {'Windows' if IS_WINDOWS else 'Non-Windows'}", "neutral"), unsafe_allow_html=True)
    if not WINDOWS_COM_AVAILABLE and COM_IMPORT_ERROR:
        st.caption(f"Detail: {COM_IMPORT_ERROR}")
    st.markdown("</div>", unsafe_allow_html=True)


elif st.session_state.active_nav == "logs":
    st.markdown('<div class="card"><div class="card-title">🖥️ System Journal</div>'
                '<div class="card-subtle">Shows the exact log lines produced by the last Dashboard pipeline run.</div>',
                unsafe_allow_html=True)
    _last_logs = st.session_state.get("last_run_log_lines")
    if _last_logs:
        st.markdown(f'<div class="terminal">{"<br>".join(_last_logs)}</div>', unsafe_allow_html=True)
    else:
        st.info("No pipeline runs yet this session. Execute the pipeline from the **Dashboard** page to populate the journal.")
    st.markdown("</div>", unsafe_allow_html=True)
